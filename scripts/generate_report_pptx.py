#!/usr/bin/env python3
"""
Generate a PowerPoint report summarizing the Abaqat al-Anwar research project.

Usage:
  uv run --with python-pptx python scripts/generate_report_pptx.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO_ROOT = Path(__file__).resolve().parent.parent
OUTPUT = REPO_ROOT / "docs" / "abaqat-research-report.pptx"

# Color palette
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_GOLD = RGBColor(0xE8, 0xB4, 0x4A)
ACCENT_TEAL = RGBColor(0x4A, 0xC4, 0xC4)
ACCENT_RED = RGBColor(0xE8, 0x5D, 0x5D)
ACCENT_GREEN = RGBColor(0x5D, 0xE8, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY = RGBColor(0x33, 0x33, 0x44)
VERY_LIGHT = RGBColor(0xF5, 0xF5, 0xF0)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=WHITE, bullet_color=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_table(slide, left, top, width, rows_data, col_widths=None, header_color=ACCENT_GOLD):
    rows = len(rows_data)
    cols = len(rows_data[0]) if rows_data else 0
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(0.3 * rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_GRAY

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xE8)

    return table_shape


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── SLIDE 1: Title ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 1, 1.5, 11, 1.5,
                 "Abaqat al-Anwar", font_size=48, bold=True, color=ACCENT_GOLD,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.0, 11, 1,
                 "عبقات الانوار في امامة الائمة الاطهار", font_size=28, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.2, 11, 0.8,
                 "Digital Research & Data Model Project", font_size=22, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 5.5, 11, 0.5,
                 "Research Report — March 2026", font_size=14, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)

    # ── SLIDE 2: What is Abaqat ──────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.4, 11, 0.8,
                 "What is Abaqat al-Anwar?", font_size=32, bold=True, color=ACCENT_GOLD)
    add_bullet_list(slide, 0.8, 1.4, 5.5, 5, [
        "23-volume systematic refutation (ردّ) by Mir Hamid Husain (d. 1306 AH / 1888 CE)",
        "Responds to Chapter 7 (باب الإمامة) of Tuhfat al-Ithna Ashariyya by Shah Abdul Aziz Dehlavi",
        "NOT a hadith collection — it is a scholarly polemic (مناظره)",
        "Defends key Shia hadiths: Ghadir, Thaqalayn, Manzilah, Safinah, Tayr, and 7 others",
        "Each volume takes ONE hadith and proves it through Sunni sources",
        "Cites 197+ unique Sunni books per volume, with full isnads",
        "Written in Farsi with embedded Arabic source quotes",
    ], font_size=14)
    add_text_box(slide, 7, 1.4, 5, 5,
                 "Structure\n\n"
                 "المنهج الأول — Quranic proofs (unpublished)\n\n"
                 "المنهج الثاني — Hadith proofs (23 volumes):\n"
                 "  Vols 1-10:  Hadith al-Ghadir\n"
                 "  Vol 11:       Hadith al-Manzilah\n"
                 "  Vol 12:       Hadith al-Wilayah\n"
                 "  Vol 13:       Hadith al-Tayr\n"
                 "  Vols 14-15: Madinat al-Ilm\n"
                 "  Vol 16:       Hadith al-Tashbih\n"
                 "  Vol 17:       Hadith al-Nur\n"
                 "  Vols 18-22: Hadith al-Thaqalayn\n"
                 "  Vol 23:       Hadith al-Safinah",
                 font_size=13, color=LIGHT_GRAY)

    # ── SLIDE 3: The Scholarly Gap ──────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.4, 11, 0.8,
                 "Dehlavi vs Mir Hamid Husain — Methodology Comparison",
                 font_size=28, bold=True, color=ACCENT_GOLD)
    add_text_box(slide, 0.8, 1.2, 11, 0.5,
                 "Based on analysis of Volume 23 (Hadith al-Safinah)", font_size=14, color=MID_GRAY)

    add_table(slide, 0.8, 1.9, 11.5, [
        ["Metric", "Shah Abdul Aziz Dehlavi", "Mir Hamid Husain"],
        ["Source books named", "0", "197 unique Sunni books"],
        ["Bare hadith quotes (no source)", "4", "0"],
        ["Isnads provided", "No", "Yes — all 40 narrator citations"],
        ["Named hadith collections", "No", "Yes — per tradition"],
        ["Citation density per argument", "0.25", "High (40 narrators, 54 scholars)"],
        ["Grading claims verified correct", "2 of 4 (50%)", "N/A (relies on named sources)"],
        ["Grading claims overstated", "2 of 4", "0"],
        ["Addresses strongest counter-argument", "No", "Yes — by section"],
        ["Time span of sources", "N/A", "204-427 AH (223 years)"],
        ["Overall scholarly grade", "D", "A"],
    ], col_widths=[3.5, 4.0, 4.0])

    # ── SLIDE 4: Dehlavi's Strategy ─────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.4, 11, 0.8,
                 "Dehlavi's Argumentation Strategy (Vol 23)",
                 font_size=28, bold=True, color=ACCENT_GOLD)
    add_bullet_list(slide, 0.8, 1.4, 5.5, 5.5, [
        "ACCEPTS the authenticity of Hadith al-Thaqalayn and Hadith al-Safinah",
        "Strategy: meaning minimization — reinterprets rather than denies",
        "",
        "Argument 1: Salvific importance ≠ political Imamate",
        "Argument 2: Counter-hadith on Rashidun Caliphs' Sunnah",
        "Argument 3: Reductio — if Safinah proves Imamate, all relatives would be Imams",
        "Argument 4: Internal Shia contradiction (Shura hadith limits Ahl al-Bayt)",
        "Argument 5: 'Following' ≠ 'political leadership'",
        "",
        "Critical weakness: 0 source books named, 50% of grading claims unverifiable",
    ], font_size=14)
    add_text_box(slide, 7, 1.4, 5, 5.5,
                 "Mir Hamid Husain's Response\n\n"
                 "• Stacks 40 Sunni narrators to prove tawatur\n"
                 "  (mass-transmission) — making denial impossible\n\n"
                 "• 11 دلالت proofs (وجوه) showing the hadith\n"
                 "  necessarily implies Imamate\n\n"
                 "• 5 biographical dossiers from rijal books\n"
                 "  evaluating contested narrators\n\n"
                 "• 10 تنبيهات (authorial notices) correcting\n"
                 "  common misconceptions\n\n"
                 "• Traces Dehlavi's arguments back to\n"
                 "  Ibn Taymiyyah — shows copying, not originality",
                 font_size=13, color=LIGHT_GRAY)

    # ── SLIDE 5: Data Model ──────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.4, 11, 0.8,
                 "Data Model — 5 Schemas for 30+ Book Types",
                 font_size=28, bold=True, color=ACCENT_GOLD)

    add_table(slide, 0.8, 1.4, 11.5, [
        ["Schema", "What it covers", "Books", "Key insight"],
        ["A: ScholarEntry", "All person-centric books", "24 books",
         "One scholar, one ID — book_type field drives which fields are populated"],
        ["B: HadithEntry", "All hadith collections", "9+ collections",
         "Isnad chain links reference scholar_id for narrator lookup"],
        ["C: ReferenceEntry", "Bibliography + Genealogical", "2 books",
         "C1 indexes books, C2 indexes nisbas"],
        ["D: Polemical", "Abaqat + Tuhfat", "2 books",
         "13 entity types in abaqat-data-model-v2.md"],
        ["E: Commentary", "Sharh/Hashiyah works", "6+ books",
         "Scholarly opinions on hadiths — cited in 15/23 volumes"],
    ], col_widths=[2.0, 3.0, 2.0, 4.5])

    add_text_box(slide, 0.8, 4.5, 11, 0.5,
                 "JarhGrade System — Standard 6+6 from al-Suyuti's Tadrib al-Rawi",
                 font_size=16, bold=True, color=ACCENT_TEAL)

    add_table(slide, 0.8, 5.2, 11.5, [
        ["Ta'dil (positive)", "Grade", "Jarh (negative)", "Grade"],
        ["أوثق الناس (superlative)", "tadil_1", "ليّن الحديث (mild)", "jarh_1"],
        ["ثقة ثبت (reinforced)", "tadil_2", "ضعيف (weak)", "jarh_2"],
        ["ثقة (single attribute)", "tadil_3", "ضعيف جداً (very weak)", "jarh_3"],
        ["صدوق (truthful)", "tadil_4", "متروك (abandoned)", "jarh_4"],
        ["صالح الحديث (acceptable)", "tadil_5", "كذّاب (liar)", "jarh_5"],
        ["مقبول (borderline)", "tadil_6", "أكذب الناس (worst)", "jarh_6"],
    ], col_widths=[3.5, 1.5, 3.5, 1.5], header_color=ACCENT_TEAL)

    # ── SLIDE 6: Digital Library ─────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.4, 11, 0.8,
                 "Digital Library — 11 Tier 1 Rijal Books",
                 font_size=28, bold=True, color=ACCENT_GOLD)
    add_text_box(slide, 0.8, 1.1, 11, 0.5,
                 "28,690 pages | 111 PDFs | 1.16 GB — All readable, searchable",
                 font_size=16, color=ACCENT_TEAL)

    add_table(slide, 0.8, 1.8, 11.5, [
        ["#", "Book", "Author", "Pages", "Verified"],
        ["1", "Tahdhib al-Tahdhib", "Ibn Hajar (d. 852)", "5,764", "✓ Ja'far al-Sadiq"],
        ["2", "Wafayat al-A'yan", "Ibn Khallikan (d. 681)", "4,087", "✓ al-Mubarrad"],
        ["3", "Mizan al-I'tidal", "al-Dhahabi (d. 748)", "2,809", "✓ Musa al-Kazim"],
        ["4", "Tadhkirat al-Huffaz", "al-Dhahabi (d. 748)", "1,760", "✓ al-Tabari/Ghadir"],
        ["5", "Mir'at al-Janan", "al-Yafi'i (d. 768)", "1,429", "✓ Shaykh al-Mufid"],
        ["6", "Al-'Ibar", "al-Dhahabi (d. 748)", "1,775", "✓ Shaykh al-Mufid"],
        ["7", "Tabaqat al-Huffaz", "al-Suyuti (d. 911)", "607", "✓ Ibn Manda"],
        ["8", "Siyar A'lam al-Nubala", "al-Dhahabi (d. 748)", "4,683", "Downloaded"],
        ["9", "Tabaqat al-Shafi'iyya", "al-Subki (d. 771)", "5,776", "Downloaded"],
        ["10", "Tarikh Baghdad", "al-Khatib (d. 463)", "~9,000", "Downloaded"],
        ["11", "Tahdhib al-Kamal", "al-Mizzi (d. 742)", "~8,000", "Downloaded"],
    ], col_widths=[0.5, 3.0, 2.5, 1.5, 2.5])

    add_text_box(slide, 0.8, 6.2, 11, 0.5,
                 "7 citations verified verbatim against source PDFs — 0 discrepancies found",
                 font_size=16, bold=True, color=ACCENT_GREEN)

    # ── SLIDE 7: Vol 23 Analysis ─────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.4, 11, 0.8,
                 "Volume 23 Deep Analysis — Hadith al-Safinah",
                 font_size=28, bold=True, color=ACCENT_GOLD)

    add_table(slide, 0.8, 1.4, 5.5, [
        ["Metric", "Count"],
        ["OCR text lines", "17,690"],
        ["Structural sections", "74"],
        ["Narrator citations (tawatur)", "40"],
        ["Dalalah proofs (وجوه)", "11"],
        ["Biographical dossiers", "5"],
        ["Tanbih notices", "10"],
        ["Total citations extracted", "518"],
        ["Unique books cited", "358"],
        ["Unique scholars", "54"],
    ], col_widths=[3.0, 2.0])

    add_text_box(slide, 7, 1.4, 5.5, 0.5,
                 "Citation Context Distribution", font_size=16, bold=True, color=ACCENT_TEAL)

    add_table(slide, 7, 2.1, 5.5, [
        ["Context", "Count", "%"],
        ["HADITH_SOURCE", "218", "42%"],
        ["PRIMARY_REBUTTAL", "117", "23%"],
        ["AUTHOR_TAWTHIQ", "77", "15%"],
        ["AUTHOR_TAJRIH", "32", "6%"],
        ["REBUTTAL_LINEAGE", "28", "5%"],
        ["CROSS_REF", "14", "3%"],
        ["Other (5 types)", "32", "6%"],
    ], col_widths=[2.5, 1.0, 1.0], header_color=ACCENT_TEAL)

    add_text_box(slide, 7, 5.5, 5.5, 1.5,
                 "Key finding: HADITH_SOURCE (42%) dominates —\n"
                 "Vol 23 focuses on proving the hadith's existence\n"
                 "in Sunni texts, not evaluating narrators.\n\n"
                 "This is the Sanad-First pattern (Pattern A):\n"
                 "Tawatur Defense → Dalalah → Dehlavi Response",
                 font_size=12, color=LIGHT_GRAY)

    # ── SLIDE 8: Hawramani + Scholar ID ──────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.4, 11, 0.8,
                 "Cross-Reference Infrastructure",
                 font_size=28, bold=True, color=ACCENT_GOLD)

    add_bullet_list(slide, 0.8, 1.5, 5.5, 3, [
        "Hawramani.com narrator index downloaded locally",
        "100,915 narrator entries from 66 source books",
        "24.5 MB — one-time download, never hit API again",
        "Enables instant name → source book lookup",
        "WordPress REST API: /wp-json/wp/v2/posts/{id}",
    ], font_size=14)

    add_text_box(slide, 7, 1.5, 5, 0.5,
                 "Scholar ID Strategy", font_size=18, bold=True, color=ACCENT_TEAL)
    add_bullet_list(slide, 7, 2.2, 5, 2.5, [
        "Primary: tahdhib-{vol}-{entry} for Tahdhib narrators",
        "Fallback: {death_ah}-{name} for non-Tahdhib scholars",
        "External: hawramani_id as cross-reference",
        "Resolution: death_ah + name fragments → match → flag ambiguous",
    ], font_size=13)

    add_text_box(slide, 0.8, 5.0, 11, 0.5,
                 "Top 5 Hawramani Source Books by Entry Count",
                 font_size=16, bold=True, color=ACCENT_TEAL)

    add_table(slide, 0.8, 5.6, 11.5, [
        ["Source Book", "Author", "Entries"],
        ["al-Jarh wa al-Ta'dil", "Ibn Abi Hatim (d. 938 CE)", "17,284"],
        ["al-Thiqat", "Ibn Hibban (d. 965 CE)", "16,018"],
        ["al-Tarikh al-Kabir", "al-Bukhari (d. 870 CE)", "13,266"],
        ["Tarikh Baghdad", "al-Khatib al-Baghdadi (d. 1071 CE)", "7,749"],
        ["Usd al-Ghaba", "Ibn al-Athir (d. 1233 CE)", "7,491"],
    ], col_widths=[4.0, 4.5, 2.0])

    # ── SLIDE 9: Next Steps ──────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.4, 11, 0.8,
                 "Next Steps", font_size=32, bold=True, color=ACCENT_GOLD)

    add_text_box(slide, 0.8, 1.5, 5.5, 0.5,
                 "Immediate", font_size=20, bold=True, color=ACCENT_GREEN)
    add_bullet_list(slide, 0.8, 2.1, 5.5, 2.5, [
        "Complete Vol 23 book classification (389 books → 6 schema types)",
        "Index Tahdhib al-Tahdhib → ScholarEntry JSON (~8,000 entries)",
        "Index Mizan al-I'tidal → ScholarEntry JSON (~11,000 entries)",
        "Build PDF-to-JSON pipeline using Gemini (not Claude — cost)",
    ], font_size=14)

    add_text_box(slide, 0.8, 4.2, 5.5, 0.5,
                 "Medium Term", font_size=20, bold=True, color=ACCENT_TEAL)
    add_bullet_list(slide, 0.8, 4.8, 5.5, 2.5, [
        "Index all 11 Tier 1 books (28,690 pages → structured JSON)",
        "Download and index 16 Tier 2 books",
        "Extend Vol 23 analysis to all 23 volumes",
        "Build cross-reference engine: scholar_id → all books → all hadiths",
    ], font_size=14)

    add_text_box(slide, 7, 1.5, 5, 0.5,
                 "Long Term Vision", font_size=20, bold=True, color=ACCENT_GOLD)
    add_bullet_list(slide, 7, 2.1, 5, 4.5, [
        "Complete digital index of Abaqat al-Anwar",
        "Every citation verified against source book",
        "Every narrator cross-referenced across 30+ books",
        "Every hadith chain traceable from companion to compiler",
        "Searchable: 'Show me all narrators Dehlavi attacks",
        "  and what Mir Hamid Husain says about each one'",
        "",
        "The first comprehensive digital apparatus",
        "for the most detailed work in Sunni-Shia",
        "hadith scholarship ever written.",
    ], font_size=14)

    # Save
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"  9 slides")


if __name__ == "__main__":
    build_presentation()
